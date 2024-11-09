import streamlit as st
from pdf2docx import Converter
from fpdf import FPDF
from docx import Document
from pptx import Presentation
import os

def convert_pdf_to_word():
    uploaded_file = st.file_uploader("Upload PDF", type="pdf")
    
    if uploaded_file is not None:
        with open("uploaded_file.pdf", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # Convert PDF to Word using pdf2docx
        converter = Converter("uploaded_file.pdf")
        converter.convert("output.docx", start=0, end=None)
        
        st.success("Conversion Complete!")
        with open("output.docx", "rb") as f:
            st.download_button("Download Word File", f, file_name="output.docx")

def convert_word_to_pdf():
    uploaded_file = st.file_uploader("Upload Word", type="docx")
    
    if uploaded_file is not None:
        with open("uploaded_file.docx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        doc = Document("uploaded_file.docx")
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        for para in doc.paragraphs:
            pdf.multi_cell(0, 10, para.text)
        
        pdf.output("output.pdf")
        
        st.success("Conversion Complete!")
        with open("output.pdf", "rb") as f:
            st.download_button("Download PDF File", f, file_name="output.pdf")

def convert_word_to_ppt():
    uploaded_file = st.file_uploader("Upload Word", type="docx")
    
    if uploaded_file is not None:
        with open("uploaded_file.docx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        doc = Document("uploaded_file.docx")
        prs = Presentation()
        
        for para in doc.paragraphs:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content Layout
            title = slide.shapes.title
            title.text = para.text
        
        prs.save("output.pptx")
        
        st.success("Conversion Complete!")
        with open("output.pptx", "rb") as f:
            st.download_button("Download PowerPoint File", f, file_name="output.pptx")

def convert_ppt_to_word():
    uploaded_file = st.file_uploader("Upload PowerPoint", type="pptx")
    
    if uploaded_file is not None:
        with open("uploaded_file.pptx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        ppt = Presentation("uploaded_file.pptx")
        doc = Document()
        
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    doc.add_paragraph(shape.text)
        
        doc.save("output.docx")
        
        st.success("Conversion Complete!")
        with open("output.docx", "rb") as f:
            st.download_button("Download Word File", f, file_name="output.docx")

def convert_ppt_to_pdf():
    uploaded_file = st.file_uploader("Upload PowerPoint", type="pptx")
    
    if uploaded_file is not None:
        with open("uploaded_file.pptx", "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # Using LibreOffice to convert PowerPoint to PDF (needs to be installed on the server)
        os.system("libreoffice --headless --convert-to pdf uploaded_file.pptx")
        
        st.success("Conversion Complete!")
        with open("uploaded_file.pdf", "rb") as f:
            st.download_button("Download PDF File", f, file_name="output.pdf")

def main():
    st.title("File Converter Application")
    
    # Select the conversion type
    option = st.radio("Choose Conversion Type", [
        "PDF to Word", 
        "Word to PDF", 
        "Word to PowerPoint", 
        "PowerPoint to Word", 
        "PowerPoint to PDF"
    ])

    if option == "PDF to Word":
        convert_pdf_to_word()
    elif option == "Word to PDF":
        convert_word_to_pdf()
    elif option == "Word to PowerPoint":
        convert_word_to_ppt()
    elif option == "PowerPoint to Word":
        convert_ppt_to_word()
    elif option == "PowerPoint to PDF":
        convert_ppt_to_pdf()

if __name__ == "__main__":
    main()